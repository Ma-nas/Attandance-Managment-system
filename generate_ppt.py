import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.text = bullets[0]
    for bullet in bullets[1:]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0

def add_image_slide(prs, title, img_path):
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only
    slide.shapes.title.text = title
    if os.path.exists(img_path):
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        slide.shapes.add_picture(img_path, left, top, width=width)

prs = Presentation()

add_title_slide(prs, 
    "Attendance Management System", 
    "A Modern, Role-Based Web Application\n\nBy [Your Name] / [USN]")

add_bullet_slide(prs, 
    "1. Problem Statement", 
    [
        "Traditional attendance methods are manual and paper-based.",
        "High risk of proxy attendance and human error.",
        "Time-consuming administrative consolidation.",
        "Lack of real-time visibility for students and administration."
    ])

add_bullet_slide(prs, 
    "2. Proposed Solution", 
    [
        "A centralized, cloud-ready web application for real-time tracking.",
        "Strict Role-Based Access Control (RBAC) for Admins, Trainers, and Students.",
        "Interactive dashboard for rapid attendance marking by trainers.",
        "Real-time analytics and transparency for students."
    ])

add_bullet_slide(prs, 
    "3. System Architecture & Tech Stack", 
    [
        "Frontend: React (Vite), TypeScript, Tailwind CSS, TanStack Query.",
        "Backend: Java 21, Spring Boot (REST APIs), Spring Data JPA.",
        "Database: PostgreSQL for reliable, relational data storage.",
        "Security: Spring Security with JWT token-based authentication.",
        "Architecture: Client-Server with isolated business logic and presentation layers."
    ])

add_image_slide(prs, "System Interface: Login Portal", "screenshots/login_page.png")
add_image_slide(prs, "System Interface: Admin Dashboard", "screenshots/admin_dashboard.png")
add_image_slide(prs, "System Interface: Trainer Interface", "screenshots/trainer_dashboard.png")

add_bullet_slide(prs, 
    "4. Conclusion & Future Scope", 
    [
        "Conclusion:",
        "- Successfully digitized the attendance workflow with a sub-second response rate.",
        "- Reduced administrative overhead significantly.",
        "Future Enhancements:",
        "- Integration with Biometric or RFID scanners.",
        "- Automated warning emails to parents for low attendance.",
        "- Machine Learning analytics to predict student dropout rates."
    ])

add_title_slide(prs, "Thank You", "Questions?")

prs.save('Attendance_Management_System_Presentation.pptx')
print("PPT generated successfully.")
